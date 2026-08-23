# -*- coding: utf-8 -*-
import sys
from pathlib import Path
import fitz

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

def pdf_text_and_scanflag(path: Path, scan_sample_pages: int = 200) -> tuple[str, bool, int]:
    try:
        doc = fitz.open(str(path))
        n_pages = doc.page_count
        text_parts = []
        empty_pages = 0
        empty_pages_sample = 0
        chars_sample = 0

        # Извлекаем текст со всех страниц
        for i in range(n_pages):
            page_text = doc.load_page(i).get_text("text")
            stripped_text = page_text.strip()
            if stripped_text:
                text_parts.append(stripped_text)
                if i < scan_sample_pages:
                    chars_sample += len(stripped_text)
            else:
                empty_pages += 1
                if i < scan_sample_pages:
                    empty_pages_sample += 1

        joined_text = "\n\n".join(text_parts)
        total_chars = len(joined_text)
        pages_examined = n_pages
        is_scan = False

        # Эвристика скана по выборке первых scan_sample_pages страниц
        if n_pages > 0 and total_chars == 0:
            is_scan = True
        elif pages_examined > 0:
            pages_examined_sample = min(n_pages, scan_sample_pages)
            if pages_examined_sample > 0:
                avg_chars_per_page_sample = chars_sample / pages_examined_sample if pages_examined_sample > 0 else 0
                empty_ratio_sample = empty_pages_sample / pages_examined_sample
                if avg_chars_per_page_sample < 50 or empty_ratio_sample > 0.7:
                    is_scan = True

        doc.close()
        return (joined_text, is_scan, n_pages)
    except Exception:
        return ("", True, 0)

def pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides_text = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_text.append(paragraph.text)
            if slide.notes_slide and slide.notes_slide.notes_text_frame:
                for paragraph in slide.notes_slide.notes_text_frame.paragraphs:
                    if paragraph.text.strip():
                        slide_text.append(f"Notes: {paragraph.text}")
            slides_text.append("\n".join(slide_text))
        return "\n\n---\n\n".join(slides_text)
    except ImportError:
        return ""
    except Exception:
        return ""

def csv_text(path: Path, max_lines: int = 3000) -> str:
    try:
        with open(path, "rb") as f:
            content = f.read()
        if content.startswith(b"\xef\xbb\xbf"):
            content = content[3:]
        text = content.decode("utf-8", errors="replace")
        lines = text.splitlines()[:max_lines]
        return "\n".join(lines)
    except Exception:
        return ""

def sheet_text(path: Path, max_rows: int = 5000) -> str:
    try:
        ext = path.suffix.lower()
        parts = []
        if ext == ".xls":
            import xlrd
            book = xlrd.open_workbook(str(path))
            for sheet in book.sheets():
                parts.append(f"## Sheet: {sheet.name}")
                for r in range(min(sheet.nrows, max_rows)):
                    cells = [str(sheet.cell_value(r, c)) for c in range(sheet.ncols)]
                    non_empty_cells = [c.strip() for c in cells if c is not None and c.strip()]
                    if non_empty_cells:
                        parts.append("\t".join(non_empty_cells))
        else:
            import openpyxl
            wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
            for ws in wb.worksheets:
                parts.append(f"## Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None and str(c).strip()]
                    if cells:
                        parts.append("\t".join(cells))
            wb.close()
        return "\n".join(parts)
    except Exception:
        return ""
